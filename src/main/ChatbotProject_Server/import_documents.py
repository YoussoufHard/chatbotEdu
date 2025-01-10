import os
import mysql.connector
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Configuration de la base de données
db_config = {
    "host": "localhost",  # Hôte de ta base de données
    "user": "root",       # Utilisateur MySQL
    "password": "",       # Mot de passe MySQL (à remplacer par ton propre mot de passe)
    "database": "chatbot_edu"  # Nom de ta base de données
}

# Chemin vers le dossier contenant les documents
directory = r'C:\Users\Lenovo\Documents\II-BDCC ENSET\BIG DATA S3\S3_2024-2025\POO JAVA\ChatbotProject_server\ENSET docs'

# Connexion à la base de données MySQL
def connect_to_db():
    return mysql.connector.connect(**db_config)

# Fonction pour lire le contenu d'un PDF
def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du PDF : {str(e)}"

#fonction pour lire le contenu d'un docx
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du DOCX : {str(e)}"
    
# Fonction pour lire le contenu d'un PPTX
def extract_text_from_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du PPTX : {str(e)}"

# Fonction pour vérifier si le document existe déjà dans la base de données
def document_exists(cursor, title):
    query = "SELECT COUNT(*) FROM documents WHERE title = %s"
    cursor.execute(query, (title,))
    result = cursor.fetchone()
    return result[0] > 0  # Si le compte est supérieur à 0, le document existe

# Fonction pour insérer un document dans la base de données
def insert_document(cursor, title, content, file_path):
    query = "INSERT INTO documents (title, content, file_path) VALUES (%s, %s, %s)"
    cursor.execute(query, (title, content, file_path))

# Fonction pour parcourir les fichiers PDF dans un dossier et ses sous-dossiers
def read_documents_recursively(directory, cursor):
    for root, _, files in os.walk(directory):
        for filename in files:
            file_path = os.path.join(root, filename)
            if filename.endswith(".pdf"):
                # Vérifier si le document existe déjà
                if not document_exists(cursor, filename):
                    # Lire le contenu du fichier PDF
                    content = extract_text_from_pdf(file_path)
                    insert_document(cursor, filename, content, file_path)
                    print(f"Document '{filename}' ajouté depuis '{root}'.")
                else:
                    print(f"Le document '{filename}' existe déjà dans la base de données.")
            elif filename.endswith(".docx"):
                content = extract_text_from_docx(file_path)
                insert_document(cursor, filename, content, file_path)
                print(f"Document Word '{filename}' ajouté.")
            elif filename.endswith(".pptx"):
                content = extract_text_from_pptx(file_path)
                insert_document(cursor, filename, content, file_path)
                print(f"Document PowerPoint '{filename}' ajouté.")


# Exécution principale
def main():
    try:
        db = connect_to_db()
        cursor = db.cursor()

        # Lire et insérer les documents du dossier et ses sous-dossiers
        read_documents_recursively(directory, cursor)

        # Commit des changements
        db.commit()
        print("Documents ajoutés avec succès à la base de données.")
    except Exception as e:
        print(f"Erreur lors du traitement : {str(e)}")
    finally:
        cursor.close()
        db.close()

if __name__ == "__main__":
    main()
